"""Turning a file on disk into text, shared by `doc` and `fetch`.

`ghostweb` reads the network; this reads `/workspace`. They share the PDF path
and nothing else, which is the whole reason this module exists: a PDF that
arrives over HTTP and the same PDF sitting in the workspace must not extract
differently, and before this there was one implementation reachable only from
the fetch path.

The case that shaped it: a file someone attached to a message. That file lands
in `/workspace/uploads/` and reaches the model as a path — for a PDF, a
spreadsheet or a scan, the path is *all* it gets, because none of those can be
put in a prompt as bytes. So the failure mode that matters is not a bad
extraction, it is an extraction that returns nothing and does not say why.
Every function here answers with either text or a sentence naming the next move.

Two decisions worth stating:

  * **The extension chooses the reader; the bytes decide whether it worked.**
    Dispatching on the name alone is how a `.pdf` that is really HTML becomes
    mojibake. Every branch checks its own output and falls through when it comes
    back empty — which is also how a text-layer-free PDF reaches OCR.

  * **OCR is opt-out, not opt-in.** A scanned PDF is indistinguishable from a
    broken one to a model holding an empty result, and "run it again with
    --ocr" is a turn spent learning something this module already knew. It runs
    automatically when the text layer is empty, says that it did, and can be
    turned off.
"""

from __future__ import annotations

import csv
import io
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path

# Long enough for a 200-page standard, short enough that one pathological file
# cannot eat the turn's whole time budget.
PDF_TIMEOUT_S = 120

# OCR is far slower than a text-layer read — seconds per page, not milliseconds.
# The page cap is what keeps "OCR this 400-page scan" from being a hang: past it
# the result says how many pages were read and how many were not.
OCR_TIMEOUT_S = 300
OCR_MAX_PAGES = 40

# The resolution `pdftoppm` renders at before OCR. 300 DPI is the figure
# tesseract's own documentation asks for; 150 halves the time and loses small
# type, which on a scanned invoice is the part someone needed.
OCR_DPI = 300

# What `doc` prints before falling back to "this is a binary I cannot read".
TEXT_SNIFF_BYTES = 8192


class DocError(Exception):
    """Extraction failed in a way worth reporting verbatim."""


@dataclass(frozen=True)
class Extracted:
    """One file, read. `note` is shown even when `text` is non-empty."""

    text: str
    kind: str
    note: str = ""


def have(command: str) -> bool:
    """Whether a helper binary is in this image. Keeps the error specific."""
    return shutil.which(command) is not None


def _run(argv: list[str], timeout: int, stdin: bytes | None = None) -> bytes:
    try:
        done = subprocess.run(
            argv, input=stdin, capture_output=True, timeout=timeout, check=False
        )
    except subprocess.TimeoutExpired as error:
        raise DocError(f"{argv[0]}: gave up after {timeout}s") from error
    except (OSError, subprocess.SubprocessError) as error:
        raise DocError(f"{argv[0]}: {error}") from error
    if done.returncode != 0:
        detail = done.stderr.decode("utf-8", "replace").strip()
        raise DocError(f"{argv[0]}: {detail or f'exit {done.returncode}'}")
    return done.stdout


def pdf_text(source: Path | bytes) -> str:
    """A PDF as text, layout preserved.

    `-layout` because a two-column paper interleaved line by line is worse than
    no text at all. Returns empty for a PDF with no text layer, which is the
    signal the caller turns into OCR rather than an error — an empty string here
    means "scanned", not "failed".
    """
    if isinstance(source, bytes):
        return _run(
            ["pdftotext", "-layout", "-nopgbrk", "-", "-"], PDF_TIMEOUT_S, stdin=source
        ).decode("utf-8", "replace").strip()
    return _run(
        ["pdftotext", "-layout", "-nopgbrk", str(source), "-"], PDF_TIMEOUT_S
    ).decode("utf-8", "replace").strip()


def pdf_pages(path: Path) -> int:
    """Page count, or 0 when `pdfinfo` will not say. Only used for reporting."""
    try:
        info = _run(["pdfinfo", str(path)], 30).decode("utf-8", "replace")
    except DocError:
        return 0
    for line in info.splitlines():
        if line.startswith("Pages:"):
            digits = line.split(":", 1)[1].strip()
            return int(digits) if digits.isdigit() else 0
    return 0


def ocr_image(path: Path, lang: str) -> str:
    """One image through tesseract."""
    if not have("tesseract"):
        raise DocError("tesseract is not installed in this toolbox")
    return _run(
        ["tesseract", str(path), "stdout", "-l", lang], OCR_TIMEOUT_S
    ).decode("utf-8", "replace").strip()


def ocr_pdf(path: Path, lang: str, max_pages: int = OCR_MAX_PAGES) -> Extracted:
    """A scanned PDF, rendered to images and read.

    Rendered with `pdftoppm` rather than `pdfimages`: a scanned page is often
    several stacked images plus a mask, and extracting them individually gives
    tesseract fragments in the wrong order. Rendering gives it the page as a
    reader sees it.
    """
    if not have("pdftoppm"):
        raise DocError("pdftoppm is not installed in this toolbox")

    total = pdf_pages(path)
    limit = min(total, max_pages) if total else max_pages

    # One page at a time, deleted as soon as it is read.
    #
    # `/tmp` is a 256 MB tmpfs — see `security.tmpfs` in the manifest — and a
    # 300 DPI A4 page is several megabytes of PNG, so rendering forty of them up
    # front is how OCR turns into "no space left on device" on exactly the long
    # scan someone most wanted read. Re-invoking `pdftoppm` per page re-parses
    # the file each time, which is milliseconds against the seconds tesseract
    # spends on the page it produced.
    read = 0
    chunks: list[str] = []
    with tempfile.TemporaryDirectory(prefix="ghost-ocr-") as work:
        for number in range(1, limit + 1):
            stem = Path(work) / f"page-{number}"
            _run(
                [
                    "pdftoppm",
                    "-r",
                    str(OCR_DPI),
                    "-f",
                    str(number),
                    "-l",
                    str(number),
                    "-png",
                    str(path),
                    str(stem),
                ],
                OCR_TIMEOUT_S,
            )
            rendered = sorted(Path(work).glob(f"page-{number}*.png"))
            if not rendered:
                break
            read += 1
            try:
                text = ocr_image(rendered[0], lang)
            finally:
                for image in rendered:
                    image.unlink(missing_ok=True)
            if text:
                chunks.append(f"--- page {number} ---\n{text}")

    if read == 0:
        raise DocError("pdftoppm produced no pages to read")

    note = f"No text layer, so {read} page(s) were read with OCR (tesseract, {lang})."
    if total and total > read:
        note += f" {total - read} further page(s) were not — raise --ocr-pages to include them."
    if not chunks:
        note = (
            f"No text layer, and OCR (tesseract, {lang}) found no words either. "
            "If the document is not in that language, name it with --lang."
        )
    return Extracted(text="\n\n".join(chunks), kind="pdf-ocr", note=note)


def read_pdf(path: Path, lang: str, ocr: bool, max_pages: int = OCR_MAX_PAGES) -> Extracted:
    """A PDF, by its text layer where it has one and by OCR where it does not."""
    # Named rather than left to `_run`, whose fallback is the raw errno —
    # "[Errno 2] No such file or directory: 'pdftotext'" reads as a problem with
    # the *document* and costs a turn to disprove.
    if not have("pdftotext"):
        raise DocError("pdftotext is not installed in this toolbox")
    text = pdf_text(path)
    if text:
        return Extracted(text=text, kind="pdf")
    if not ocr:
        return Extracted(
            text="",
            kind="pdf",
            note=(
                "This PDF has no text layer — it is page images. OCR is available "
                "here; re-run without --no-ocr to read it."
            ),
        )
    return ocr_pdf(path, lang, max_pages)


def read_image(path: Path, lang: str, ocr: bool) -> Extracted:
    """An image: what it is, and any words in it.

    The dimensions are printed even when there is no text, because "this is a
    3024×4032 photo" and "this file did not open" are different situations and a
    bare empty result cannot tell them apart.
    """
    size = ""
    try:
        from PIL import Image  # noqa: PLC0415 — optional, and only on this path

        with Image.open(path) as handle:
            size = f"{handle.width}×{handle.height} {handle.format or ''}".strip()
    except Exception:  # noqa: BLE001 — a broken image is a note, not a crash
        size = ""

    note = f"Image{f' · {size}' if size else ''}."
    if not ocr:
        return Extracted(text="", kind="image", note=f"{note} OCR was not run.")

    text = ocr_image(path, lang)
    if not text:
        note += (
            f" OCR (tesseract, {lang}) found no words. If it holds text in another "
            "language, name it with --lang."
        )
    return Extracted(text=text, kind="image-ocr", note=note)


def read_docx(path: Path) -> Extracted:
    """A Word document: paragraphs and tables, in document order."""
    try:
        import docx  # noqa: PLC0415
    except ImportError as error:  # pragma: no cover — depends on the image
        raise DocError("python-docx is not installed in this toolbox") from error

    document = docx.Document(str(path))
    blocks = [paragraph.text.rstrip() for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                blocks.append(" | ".join(cells))
    text = "\n".join(block for block in blocks if block).strip()
    return Extracted(text=text, kind="docx", note="" if text else "The document is empty.")


def _cell(value: object) -> str:
    """One spreadsheet cell as text.

    A date-only cell arrives as a `datetime` at midnight, and `str()` on that is
    `2026-03-01 00:00:00` — three fields of noise on every row of what is usually
    the most-read column in the sheet.
    """
    if value is None:
        return ""
    if isinstance(value, datetime):
        if (value.hour, value.minute, value.second, value.microsecond) == (0, 0, 0, 0):
            return value.date().isoformat()
        return value.isoformat(sep=" ")
    if isinstance(value, date):
        return value.isoformat()
    return str(value)


def read_xlsx(path: Path, max_rows: int = 2000) -> Extracted:
    """A spreadsheet, one CSV block per sheet.

    CSV rather than a rendered grid because a model reads it more reliably and it
    is a third of the tokens. `data_only` so a formula cell gives its last saved
    value — the formula text is almost never what was wanted, and a sheet of
    `=SUM(B2:B40)` reads as an empty sheet.
    """
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    except ImportError as error:  # pragma: no cover
        raise DocError("openpyxl is not installed in this toolbox") from error

    book = load_workbook(str(path), read_only=True, data_only=True)
    blocks: list[str] = []
    truncated = False
    for sheet in book.worksheets:
        buffer = io.StringIO()
        # The `csv` module rather than `",".join`, so a cell containing a comma
        # is quoted instead of silently splitting the row. Joining by hand and
        # stripping commas out of the data was the first version, and it made
        # "Smith, J" and "Smith  J" the same value.
        writer = csv.writer(buffer, lineterminator="\n")
        wrote = 0
        for index, row in enumerate(sheet.iter_rows(values_only=True)):
            if index >= max_rows:
                truncated = True
                break
            cells = [_cell(value) for value in row]
            while cells and cells[-1] == "":
                cells.pop()
            if cells:
                writer.writerow(cells)
                wrote += 1
        if wrote:
            blocks.append(f"--- sheet: {sheet.title} ---\n" + buffer.getvalue().rstrip())
    book.close()

    note = f"{len(blocks)} sheet(s)."
    if truncated:
        note += f" Stopped at {max_rows} rows per sheet — raise --max-rows for the rest."
    return Extracted(text="\n\n".join(blocks), kind="xlsx", note=note)


def read_pptx(path: Path) -> Extracted:
    """A deck: every slide's text, including speaker notes."""
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError as error:  # pragma: no cover
        raise DocError("python-pptx is not installed in this toolbox") from error

    deck = Presentation(str(path))
    blocks: list[str] = []
    for number, slide in enumerate(deck.slides, start=1):
        lines = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[notes] {notes}")
        if lines:
            blocks.append(f"--- slide {number} ---\n" + "\n".join(lines))
    return Extracted(
        text="\n\n".join(blocks),
        kind="pptx",
        note=f"{len(blocks)} slide(s) with text.",
    )


def looks_binary(head: bytes) -> bool:
    """Whether these bytes are not text, by the two rules that cost nothing.

    A NUL first, which is the rule `git` uses and the one the host's own file
    tools use — keeping the same answer on both sides of the container matters
    more than being clever, because a `doc` that disagrees with `read_file`
    about what is readable is worse than either rule being imperfect.

    Then a share of control bytes, because NUL alone is a coin flip on a short
    file: 200 bytes of compressed data contain one only about half the time, and
    the miss prints a screenful of mojibake. Tab, newline, carriage return and
    form feed are text; the rest of C0 is not, and no real document is a third
    control characters.

    Only C0 counts. Everything at 0x80 and above is left alone because that is
    where UTF-8 keeps its continuation bytes — treating those as suspicious
    declared `héllo` and `русский текст` binary, which is the one thing a
    toolbox shipping thirteen OCR languages must not do.

    What this deliberately does not catch: a short file of uniformly random
    bytes, which has too few C0 codes to trip the ratio and often no NUL either.
    A strict UTF-8 decode would catch it and would also declare every latin-1
    document binary, which is a worse trade for these languages. Real binaries —
    archives, images, office files, anything compressed — are dense with NULs
    well inside the first pages, so the gap is a synthetic case rather than a
    file anyone attaches.
    """
    if b"\0" in head:
        return True
    if not head:
        return False
    textish = b"\t\n\r\f\b\x1b" + bytes(range(0x20, 0x7F)) + bytes(range(0x80, 0x100))
    return len(head.translate(None, delete=textish)) * 3 > len(head)


def read_text(path: Path) -> Extracted | None:
    """The file as text, or `None` when its bytes say it is not text.

    Decided from the bytes, never the extension, which is wrong in both
    directions: `.py` and `.ts` have no MIME entry anywhere, and a `.txt` can
    hold a binary blob.
    """
    with path.open("rb") as handle:
        head = handle.read(TEXT_SNIFF_BYTES)
    if looks_binary(head):
        return None
    return Extracted(text=path.read_text("utf-8", errors="replace").strip(), kind="text")
